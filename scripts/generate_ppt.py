from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

SLIDES_MD = Path(__file__).resolve().parents[1] / 'presentations' / 'spotify_popularity_model_slides.md'
OUTPUT_PPTX = Path(__file__).resolve().parents[1] / 'presentations' / 'spotify_popularity_model_slides.pptx'


def add_bullets(text_frame, lines):
    text_frame.clear()
    if not lines:
        return
    # First line as title-like bullet
    p = text_frame.paragraphs[0]
    p.text = lines[0]
    p.level = 0
    for line in lines[1:]:
        p = text_frame.add_paragraph()
        p.text = line
        p.level = 1


def parse_sections(md_text):
    sections = []
    current_title = None
    current_bullets = []
    for raw in md_text.splitlines():
        line = raw.strip()
        if line.startswith('### '):
            if current_title is not None:
                sections.append((current_title, current_bullets))
            current_title = line.replace('### ', '').strip()
            current_bullets = []
        elif line.startswith('- '):
            current_bullets.append(line[2:])
    if current_title is not None:
        sections.append((current_title, current_bullets))
    return sections


def main():
    md = SLIDES_MD.read_text(encoding='utf-8')
    sections = parse_sections(md)

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content

    for title, bullets in sections[:3]:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        add_bullets(tf, bullets)

    prs.save(OUTPUT_PPTX)
    print(f"Saved: {OUTPUT_PPTX}")


if __name__ == '__main__':
    main()